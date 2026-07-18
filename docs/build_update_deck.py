"""Builds the 2-week project-update deck (with live screenshots) as a .pptx.

Run: backend/.venv/Scripts/python.exe docs/build_update_deck.py
Output: docs/VHN_Project_Update.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(__file__)
SHOTS = os.path.join(HERE, "screenshots")

# palette
PRIMARY = RGBColor(0x4F, 0x6D, 0xF5)
PRIMARY_LT = RGBColor(0xAB, 0xBA, 0xFB)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
TEAL = RGBColor(0x10, 0x9D, 0x76)
DARK = RGBColor(0x0E, 0x16, 0x30)
DARK_CARD = RGBColor(0x1A, 0x24, 0x46)
LIGHTBG = RGBColor(0xF4, 0xF6, 0xFB)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
MUTED_LT = RGBColor(0xAE, 0xB9, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SOFT = RGBColor(0xEE, 0xF1, 0xFE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         first=False, sb=0, sa=4, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Segoe UI"
    return p


def card(s, l, t, w, h, fill=WHITE, border=LINE, radius=0.07, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def heading(s, title, sub, ink=INK, sub_color=MUTED):
    para(tb(s, 0.7, 0.42, 12, 0.9), title, 32, ink, bold=True, first=True)
    if sub:
        para(tb(s, 0.72, 1.28, 12, 0.5), sub, 13.5, sub_color, first=True)


def phone_shot(s, name, cx, top, h):
    """Place a phone screenshot centered at cx, scaled to height h (inches)."""
    path = os.path.join(SHOTS, f"{name}.png")
    if not os.path.exists(path):
        return
    pic = s.shapes.add_picture(path, Inches(0), Inches(top), height=Inches(h))
    w = pic.width
    pic.left = Inches(cx) - int(w / 2)
    # rounded-ish framing: a subtle border behind
    return pic


def caption(s, text, cx, top, w=3.0):
    tf = tb(s, cx - w / 2, top, w, 0.5, anchor=MSO_ANCHOR.TOP)
    para(tf, text, 11.5, INK, bold=True, align=PP_ALIGN.CENTER, first=True)


# ================= SLIDE 1 — Title =================
s = slide(DARK)
sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.7), Inches(-2.2), Inches(6.5), Inches(6.5))
sh.fill.solid(); sh.fill.fore_color.rgb = DARK_CARD; sh.line.fill.background(); sh.shadow.inherit = False
para(tb(s, 0.85, 1.5, 11, 1.1), "Virtual Health Navigator", 44, WHITE, bold=True, first=True)
para(tb(s, 0.9, 2.65, 11, 0.5), "Project Update — Weeks 7–8", 21, PRIMARY_LT, first=True)
para(tb(s, 0.9, 3.35, 11, 0.4), "Live OpenBioLLM · Doctor dashboard · Live wait times · Dark mode · Responsive",
     13, MUTED_LT, first=True)
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.0), Inches(3.4), Inches(0.09))
bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background(); bar.shadow.inherit = False
para(tb(s, 0.9, 4.4, 11, 0.4), "CSIS 4495 – 071  ·  Team MAYA", 14, MUTED_LT, first=True)
para(tb(s, 0.9, 4.85, 11, 0.4), "Amish Nanda · Lovepreet Singh · Shinsuke Tomita", 13, MUTED_LT, first=True)
para(tb(s, 9.4, 6.8, 3.3, 0.4), "July 2026", 12, MUTED_LT, align=PP_ALIGN.RIGHT, first=True)

# ================= SLIDE 2 — What we shipped =================
s = slide(LIGHTBG)
heading(s, "What we shipped this update", "Six major features added since the midterm")
items = [
    ("Live OpenBioLLM AI", "Real biomedical LLM (OpenBioLLM-8B) now powers the symptom checker "
     "and chat — hosted on Hugging Face, no local GPU needed.", VIOLET, RGBColor(0xF3, 0xEF, 0xFD)),
    ("Health AI chat", "A conversational assistant behind a floating button; "
     "tap any past check to continue it in chat.", PRIMARY, SOFT),
    ("Doctor dashboard", "New doctor role — patients pick a doctor, who sees their roster "
     "and symptom-check history.", TEAL, RGBColor(0xE7, 0xF6, 0xF0)),
    ("Live ER wait times", "Real Metro Vancouver ED / urgent-care wait times via edwaittimes.ca, "
     "merged into the locator with distance sorting.", ORANGE, RGBColor(0xFD, 0xF6, 0xE9)),
    ("Dark mode", "A theme toggle persisted per-user in the database (no localStorage).",
     RGBColor(0x33, 0x41, 0x55), RGBColor(0xE9, 0xEE, 0xF6)),
    ("Responsive + redesign", "Mobile-first layout that becomes a sidebar app on tablet/desktop; "
     "a Samsung-Health-style profile.", RGBColor(0xE0, 0x79, 0x8A), RGBColor(0xFD, 0xEC, 0xEF)),
]
fl, fgap = 0.6, 0.3
fw = (12.13 - fgap) / 2
ftop, fh, fvgap = 1.75, 1.28, 0.16
for i, (title, body, accent, fill) in enumerate(items):
    col, row = i % 2, i // 2
    x = fl + col * (fw + fgap)
    y = ftop + row * (fh + fvgap)
    card(s, x, y, fw, fh, fill=fill, border=None, radius=0.08)
    tfc = tb(s, x + 0.3, y + 0.2, fw - 0.6, fh - 0.35)
    para(tfc, title, 14, accent, bold=True, first=True, sa=4)
    para(tfc, body, 10.5, MUTED, line=1.08)

# ================= SLIDE 3 — Live demo: symptom checker + AI chat =================
s = slide(LIGHTBG)
heading(s, "Live demo — AI symptom guidance", "OpenBioLLM interprets symptoms; the rule engine stays the safety floor")
phone_shot(s, "04_symptom_checker", 3.3, 1.7, 5.0)
phone_shot(s, "05_chat_openbiollm", 6.9, 1.7, 5.0)
caption(s, "Symptom checker", 3.3, 6.75)
caption(s, "Health AI chat — “Reviewed by OpenBioLLM”", 6.9, 6.75, w=4.2)
tf = tb(s, 10.2, 2.0, 2.9, 4.5)
para(tf, "How it works", 14, PRIMARY, bold=True, first=True, sa=8)
for b in ["Symptoms → OpenBioLLM-8B on Hugging Face (Featherless provider).",
          "Model may reword the guidance…", "…but the rule engine decides urgency.",
          "Emergencies bypass the model entirely.", "Cold-start falls back to rules — app never breaks."]:
    para(tf, "•  " + b, 10.5, INK, sa=7, line=1.1)

# ================= SLIDE 4 — Live demo: home, dark, profile =================
s = slide(LIGHTBG)
heading(s, "Live demo — dashboard, dark mode, profile", "A native-feeling PWA with a Samsung-Health-style profile")
phone_shot(s, "02_home_light", 2.7, 1.7, 5.0)
phone_shot(s, "03_home_dark", 6.65, 1.7, 5.0)
phone_shot(s, "08_profile_mypage", 10.6, 1.7, 5.0)
caption(s, "Home dashboard", 2.7, 6.75)
caption(s, "Dark mode (saved per user)", 6.65, 6.75, w=3.6)
caption(s, "Profile / My page", 10.6, 6.75)

# ================= SLIDE 5 — Live demo: doctor + wait times =================
s = slide(LIGHTBG)
heading(s, "Live demo — doctor dashboard & live wait times", "A multi-sided model + real BC emergency-department data")
phone_shot(s, "10_doctor_dashboard", 2.7, 1.7, 5.0)
phone_shot(s, "11_doctor_patient_history", 6.65, 1.7, 5.0)
phone_shot(s, "06_nearby_waittimes", 10.6, 1.7, 5.0)
caption(s, "Doctor: patient roster", 2.7, 6.75)
caption(s, "Doctor: patient history", 6.65, 6.75)
caption(s, "Nearby: live ER wait times", 10.6, 6.75)

# ================= SLIDE 6 — Responsive =================
s = slide(LIGHTBG)
heading(s, "Responsive — mobile-first, adapts to any screen", "The bottom tab bar becomes a sidebar on tablet / laptop")
pic = s.shapes.add_picture(os.path.join(SHOTS, "12_desktop_sidebar.png"), Inches(0), Inches(1.75), width=Inches(8.4))
pic.left = Inches(0.7)
card(s, 9.4, 1.9, 3.4, 3.6, fill=WHITE, border=LINE, radius=0.06)
tf = tb(s, 9.7, 2.1, 2.8, 3.3)
para(tf, "One codebase", 14, PRIMARY, bold=True, first=True, sa=8)
for b in ["Phone → full-screen app with bottom tabs.",
          "Tablet → roomier layout, centered content.",
          "Laptop → left sidebar + wide content card.",
          "Touch devices fill the screen; desktop shows the app card.",
          "Built with CSS grid + :has() — no JS."]:
    para(tf, "•  " + b, 10.5, INK, sa=7, line=1.1)
caption(s, "Same app on a laptop — sidebar navigation", 4.9, 6.85, w=6)

# ================= SLIDE 7 — New techniques learned =================
s = slide(DARK)
para(tb(s, 0.7, 0.45, 12, 0.9), "New techniques we learned", 32, WHITE, bold=True, first=True)
para(tb(s, 0.72, 1.3, 12, 0.5), "Skills picked up building this update", 13.5, MUTED_LT, first=True)
tech = [
    ("Hosted LLM inference", VIOLET, ["HF Inference Providers API (Featherless)",
     "Llama-3 chat templates & text-generation task", "A safety floor the model can never override"]),
    ("Full-stack auth & data", TEAL, ["DB-backed sessions + HTTP-only cookies (no localStorage)",
     "PBKDF2 password hashing", "Non-destructive ALTER TABLE migrations", "Role-based access (doctor vs patient)"]),
    ("Responsive & integration", AMBER, ["CSS grid + :has() for role-based layouts",
     "pointer/hover media queries (touch vs mouse)", "Consuming an undocumented API with caching + fallback",
     "Automated screenshots with Playwright"]),
]
cl, cgap = 0.65, 0.4
cw = (12.0 - cgap * 2) / 3
for i, (title, accent, lines) in enumerate(tech):
    x = cl + i * (cw + cgap)
    card(s, x, 1.95, cw, 4.3, fill=DARK_CARD, border=accent, radius=0.05, lw=1.5)
    para(tb(s, x + 0.28, 2.2, cw - 0.55, 0.5), title, 14.5, accent, bold=True, first=True)
    tf = tb(s, x + 0.28, 2.95, cw - 0.5, 3.1)
    for j, b in enumerate(lines):
        para(tf, "•  " + b, 11, MUTED_LT, first=(j == 0), sa=9, line=1.12)
para(tb(s, 0.7, 6.55, 12, 0.5),
     "VHN — guidance, not diagnosis. Every result routes to a real clinician.",
     12, PRIMARY_LT, italic=True, align=PP_ALIGN.CENTER, first=True)

out = os.path.join(HERE, "VHN_Project_Update.pptx")
prs.save(out)
print("Saved:", out)
