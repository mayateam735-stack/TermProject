"""Builds the refined VHN midterm deck (HealthNav indigo theme) as a .pptx.

Run:  backend/.venv/Scripts/python.exe docs/build_midterm_deck.py
Output: docs/VHN_Midterm_Refined.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ----
PRIMARY = RGBColor(0x4F, 0x6D, 0xF5)
PRIMARY_LT = RGBColor(0xAB, 0xBA, 0xFB)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
TEAL = RGBColor(0x10, 0x9D, 0x76)
DARK = RGBColor(0x0E, 0x16, 0x30)
DARK_CARD = RGBColor(0x1A, 0x24, 0x46)
LIGHTBG = RGBColor(0xF4, 0xF6, 0xFB)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
MUTED_LT = RGBColor(0xAE, 0xB9, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SOFT = RGBColor(0xEE, 0xF1, 0xFE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         first=False, sb=0, sa=4, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = "Segoe UI"
    return p


def card(s, l, t, w, h, fill=WHITE, border=LINE, radius=0.06, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def oval(s, l, t, w, h, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def heading(s, title, sub, ink=INK, sub_color=MUTED):
    para(tb(s, 0.7, 0.45, 12, 0.9), title, 33, ink, bold=True, first=True)
    para(tb(s, 0.72, 1.32, 12, 0.5), sub, 13.5, sub_color, first=True)


def card_block(s, l, t, w, h, title, lines, accent, fill=WHITE, border=LINE,
               title_size=14.5, body_size=10.5, body_color=MUTED, lw=1.0):
    card(s, l, t, w, h, fill=fill, border=border, lw=lw, radius=0.07)
    tf = tb(s, l + 0.28, t + 0.22, w - 0.56, h - 0.4)
    para(tf, title, title_size, accent, bold=True, first=True, sa=5)
    for ln in lines:
        para(tf, ln, body_size, body_color, sa=2, line=1.05)


# ================= SLIDE 1 — Title =================
s = slide(DARK)
oval(s, 9.7, -2.2, 6.5, 6.5, DARK_CARD)
oval(s, -1.6, 4.6, 3.6, 3.6, RGBColor(0x16, 0x20, 0x40))
para(tb(s, 0.8, 0.9, 11, 1.1), "Virtual Health Navigator", 46, WHITE, bold=True, first=True)
para(tb(s, 0.85, 2.05, 11, 0.5), "Midterm Progress Presentation", 21, PRIMARY_LT, first=True)
para(tb(s, 0.85, 2.72, 11, 0.4), "CSIS 4495 – 071  ·  Team MAYA", 14, MUTED_LT, first=True)
para(tb(s, 0.85, 3.18, 11, 0.4), "Amish Nanda  ·  Lovepreet Singh  ·  Shinsuke Tomita", 14, MUTED_LT, first=True)
bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(3.85), Inches(3.4), Inches(0.09))
bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background(); bar.shadow.inherit = False

pills = ["Symptom Checker", "ER Decision Flow", "Clinic Locator",
         "Medication Reminders", "Health AI Chat", "Health History"]
pl, pw, pgap, py, ph = 0.8, 1.83, 0.13, 4.55, 1.55
for i, label in enumerate(pills):
    x = pl + i * (pw + pgap)
    card(s, x, py, pw, ph, fill=DARK_CARD, border=RGBColor(0x2A, 0x36, 0x60), radius=0.08, lw=1.0)
    oval(s, x + pw / 2 - 0.22, py + 0.35, 0.44, 0.44, PRIMARY)
    tf = tb(s, x + 0.1, py + 0.95, pw - 0.2, 0.5, anchor=MSO_ANCHOR.TOP)
    para(tf, label, 10.5, RGBColor(0xC9, 0xD3, 0xEC), bold=True, align=PP_ALIGN.CENTER, first=True)
para(tb(s, 9.5, 6.85, 3.2, 0.4), "June 2026", 12, MUTED_LT, align=PP_ALIGN.RIGHT, first=True)

# ================= SLIDE 2 — Progress =================
s = slide(LIGHTBG)
heading(s, "Overall Project Progress", "9-week plan  ·  Weeks 0–6 complete")
phases = [
    ("Wk 0", "Discovery", "Scope & requirements", True),
    ("Wk 1", "UX & Wireframes", "Flows, mockups, design", True),
    ("Wk 2–4", "Core Build", "Checker, ER flow, profile", True),
    ("Wk 5–6", "Locator + Reminders", "Clinics, wait times", True),
    ("Wk 7–8", "Advanced Module", "AI summary, voice, PDF", False),
    ("Wk 9", "Test & Launch", "QA, portal, final demo", False),
]
pl, pgap = 0.6, 0.18
pw = (12.13 - pgap * 5) / 6
ptop = 2.0
for i, (wk, title, desc, done) in enumerate(phases):
    x = pl + i * (pw + pgap)
    hd = card(s, x, ptop, pw, 0.5, fill=(PRIMARY if done else RGBColor(0xCB, 0xD5, 0xE1)), border=None, radius=0.18)
    para(tb(s, x, ptop + 0.1, pw, 0.35, anchor=MSO_ANCHOR.MIDDLE), wk, 12.5,
         WHITE if done else RGBColor(0x47, 0x55, 0x69), bold=True, align=PP_ALIGN.CENTER, first=True)
    body_fill = RGBColor(0xEC, 0xFB, 0xF4) if done else WHITE
    body_border = RGBColor(0xBF, 0xE8, 0xD8) if done else LINE
    card(s, x, ptop + 0.62, pw, 1.9, fill=body_fill, border=body_border, radius=0.07)
    tf = tb(s, x + 0.14, ptop + 0.8, pw - 0.28, 1.6)
    para(tf, ("✓ " if done else "") + title, 12, INK if done else MUTED, bold=True, first=True, sa=5, line=1.0)
    para(tf, desc, 10, MUTED, sa=2, line=1.05)

# progress bar
card(s, 0.6, 4.95, 12.13, 0.4, fill=RGBColor(0xE2, 0xE8, 0xF0), border=None, radius=0.5)
card(s, 0.6, 4.95, 12.13 * 0.67, 0.4, fill=TEAL, border=None, radius=0.5)
para(tb(s, 0.6, 5.45, 12.13, 0.4), "~67% complete  (Weeks 0–6 of 9)", 13, INK, bold=True,
     align=PP_ALIGN.CENTER, first=True)

roles = [
    ("Amish Nanda", "Backend · FastAPI · PostgreSQL · LLM", VIOLET, RGBColor(0xF2, 0xEE, 0xFC)),
    ("Lovepreet Singh", "Frontend · React PWA · UI/UX", TEAL, RGBColor(0xE7, 0xF6, 0xF0)),
    ("Shinsuke Tomita", "AI/Safety · Triage logic · Privacy", ORANGE, RGBColor(0xFD, 0xF3, 0xE3)),
]
rl, rgap = 0.6, 0.25
rw = (12.13 - rgap * 2) / 3
for i, (name, role, col, fill) in enumerate(roles):
    x = rl + i * (rw + rgap)
    card(s, x, 6.05, rw, 1.05, fill=fill, border=None, radius=0.1)
    tf = tb(s, x, 6.22, rw, 0.9, anchor=MSO_ANCHOR.TOP)
    para(tf, name, 13, col, bold=True, align=PP_ALIGN.CENTER, first=True, sa=4)
    para(tf, role, 10.5, MUTED, align=PP_ALIGN.CENTER)

# ================= SLIDE 3 — Architecture =================
s = slide(LIGHTBG)
heading(s, "System Architecture & Tech Stack",
        "Three-tier architecture — React PWA · FastAPI · PostgreSQL on Neon")
tiers = [
    ("Tier 1 — Frontend", "React PWA · Vite · localhost:5173", TEAL, RGBColor(0xEC, 0xFA, 0xF4),
     ["Symptom checker UI", "Health AI chat", "Home dashboard", "Clinic locator map",
      "Medication reminders", "Auth / session"]),
    ("Tier 2 — API / Service", "Python · FastAPI · localhost:8000", VIOLET, RGBColor(0xF3, 0xEF, 0xFD),
     ["Auth router", "Triage router", "Chat router", "Symptom-checks router",
      "Clinics router", "Reminders · Patients routers"]),
    ("Tier 3 — Data + AI", "PostgreSQL (Neon) · SQLAlchemy · LLM stub", ORANGE, RGBColor(0xFD, 0xF6, 0xE9),
     ["PostgreSQL on Neon", "Server-side sessions", "Rule-based safety floor",
      "LLM stub (OpenBioLLM-8B)", "RAG integration (planned)", "BC clinic seeder"]),
]
tl, tgap = 0.55, 0.42
tw = (12.23 - tgap * 2) / 3
ttop, th = 1.75, 4.5
for i, (title, sub, accent, fill, items) in enumerate(tiers):
    x = tl + i * (tw + tgap)
    card(s, x, ttop, tw, th, fill=fill, border=accent, radius=0.05, lw=1.5)
    para(tb(s, x + 0.25, ttop + 0.2, tw - 0.5, 0.5), title, 15.5, accent, bold=True, first=True)
    para(tb(s, x + 0.25, ttop + 0.72, tw - 0.5, 0.4), sub, 10, MUTED, italic=True, first=True)
    iy, ih, gap = ttop + 1.2, 0.44, 0.1
    for it in items:
        card(s, x + 0.22, iy, tw - 0.44, ih, fill=WHITE, border=LINE, radius=0.14, lw=0.75)
        para(tb(s, x + 0.4, iy, tw - 0.7, ih, anchor=MSO_ANCHOR.MIDDLE), it, 10.5, INK, first=True)
        iy += ih + gap
    if i < 2:
        para(tb(s, x + tw + 0.04, ttop + 1.9, 0.34, 0.5, anchor=MSO_ANCHOR.MIDDLE), "→", 20, MUTED,
             bold=True, align=PP_ALIGN.CENTER, first=True)
cap = card(s, 0.55, 6.5, 12.23, 0.6, fill=RGBColor(0xE7, 0xF6, 0xF0), border=None, radius=0.18)
para(tb(s, 0.55, 6.5, 12.23, 0.6, anchor=MSO_ANCHOR.MIDDLE),
     "🛡  Rule-based safety floor always overrides the LLM — the model can never downgrade urgency",
     12.5, TEAL, bold=True, align=PP_ALIGN.CENTER, first=True)

# ================= SLIDE 4 — Current Implementation =================
s = slide(LIGHTBG)
heading(s, "Current Implementation Overview",
        "Working MVP — all core features functional end-to-end")
feats = [
    ("Auth & Sessions", ["DB-stored server-side sessions + HTTP-only cookie; PBKDF2-HMAC-SHA256",
                          "hashed passwords. Signup / login / logout / me endpoints live."], VIOLET, RGBColor(0xF3, 0xEF, 0xFD)),
    ("Symptom Checker", ["POST /api/triage — rule engine scores tags + pain level + duration",
                         "into 4 care levels with ER safety logic."], TEAL, RGBColor(0xEC, 0xFA, 0xF4)),
    ("Health AI Chat", ["POST /api/chat — conversational guidance behind a floating button;",
                        "same safety floor (emergencies flagged in chat too)."], VIOLET, RGBColor(0xF3, 0xEF, 0xFD)),
    ("Home Dashboard", ["Greeting, quick-action tiles (incl. tap-to-call 8-1-1 / 911),",
                        "live Upcoming reminders & Recent activity."], PRIMARY, SOFT),
    ("Clinic Locator", ["GET /api/clinics with ?kind=&lat=&lng= geo-filter + distance sort.",
                        "BC sample clinics seeded via app.seed."], ORANGE, RGBColor(0xFD, 0xF6, 0xE9)),
    ("Medication Reminders", ["CRUD + PATCH /{id}/taken. Per-patient and DB-persisted —",
                              "including the 'taken today' state (PostgreSQL on Neon)."], RED, RGBColor(0xFD, 0xEC, 0xEC)),
    ("Health History & Profile", ["Every check stored per patient; GET /api/patients/me/history.",
                                  "Editable profile via PATCH /api/patients/me + shareable summary."], GREEN, RGBColor(0xEA, 0xF7, 0xEF)),
    ("LLM Integration Point", ["OpenBioLLM-8B stubbed in llm.py — app runs without a model.",
                               "Rule engine acts as a permanent safety floor."], VIOLET, RGBColor(0xF3, 0xEF, 0xFD)),
]
fl, fgap = 0.6, 0.3
fw = (12.13 - fgap) / 2
ftop, fh, fvgap = 1.75, 1.0, 0.14
for i, (title, lines, accent, fill) in enumerate(feats):
    col, row = i % 2, i // 2
    x = fl + col * (fw + fgap)
    y = ftop + row * (fh + fvgap)
    card_block(s, x, y, fw, fh, title, lines, accent, fill=fill, border=None,
               title_size=13.5, body_size=9.5)
para(tb(s, 0.6, 6.62, 12.13, 0.5, anchor=MSO_ANCHOR.MIDDLE),
     "Interactive API docs at /docs  ·  Frontend at :5173  ·  Every record scoped to the signed-in patient",
     11, TEAL, bold=True, align=PP_ALIGN.CENTER, first=True)

# ================= SLIDE 5 — Challenges =================
s = slide(DARK)
para(tb(s, 0.7, 0.45, 12, 0.9), "Challenges & Issues", 33, WHITE, bold=True, first=True)
para(tb(s, 0.72, 1.32, 12, 0.5), "Known constraints and planned mitigations", 13.5, MUTED_LT, first=True)
chal = [
    ("LLM inference on student hardware", VIOLET,
     ["OpenBioLLM-8B at 8-bit needs ~8GB VRAM. Student laptops and free-tier cloud",
      "(Vercel/Heroku) can't host it. Mitigation: the stub ships the app now; a GPU",
      "instance (RunPod / local) is targeted for the advanced module."]),
    ("Safety-floor correctness", RED,
     ["Triage must never miss a life-threatening keyword; pain-level & duration scoring",
      "now layer on top. Edge cases (ambiguous, multi-symptom) need exhaustive tests.",
      "Mitigation: growing keyword list + integration tests on every change."]),
    ("Live clinic data & wait times", AMBER,
     ["BC's Health Connect Registry API needs institutional access unavailable to",
      "student projects; wait times are mocked from the seeder. Mitigation: static",
      "seed data clearly labelled 'estimated'; live integration deferred."]),
    ("Neon DB schema migrations", TEAL,
     ["SQLAlchemy create_all only adds missing tables — never columns. Adding a model",
      "field without ALTER TABLE causes UndefinedColumn / 500 errors. Mitigation:",
      "ALTER TABLE or drop-and-recreate on schema changes (Alembic planned)."]),
]
cl, cgap = 0.65, 0.4
cw = (12.0 - cgap) / 2
ctop, ch, cvgap = 1.9, 2.05, 0.3
for i, (title, accent, lines) in enumerate(chal):
    col, row = i % 2, i // 2
    x = cl + col * (cw + cgap)
    y = ctop + row * (ch + cvgap)
    card(s, x, y, cw, ch, fill=DARK_CARD, border=accent, radius=0.05, lw=1.5)
    tf = tb(s, x + 0.32, y + 0.28, cw - 0.6, ch - 0.5)
    para(tf, title, 15, accent, bold=True, first=True, sa=7)
    for ln in lines:
        para(tf, ln, 10.5, MUTED_LT, sa=2, line=1.1)

# ================= SLIDE 6 — AI Usage =================
s = slide(LIGHTBG)
heading(s, "AI Usage in Implementation",
        "Claude (Anthropic) used as a development assistant — all outputs reviewed and verified by the team")
cols = [
    ("Code scaffolding", TEAL, RGBColor(0xEC, 0xFA, 0xF4),
     ["FastAPI router stubs & SQLAlchemy models",
      "React component structure & API proxy config",
      "Seed script for BC sample clinics"]),
    ("Triage logic design", VIOLET, RGBColor(0xF3, 0xEF, 0xFD),
     ["Keyword taxonomy for ER escalation rules",
      "Safety-floor invariant — LLM can't downgrade",
      "Pain-level & duration scoring over the floor",
      "Test cases for edge-case symptoms"]),
    ("Research & documentation", ORANGE, RGBColor(0xFD, 0xF6, 0xE9),
     ["B.C. healthcare stats (400K unattached, 32-wk wait)",
      "Literature: Gilbert et al. 2020, Wallace et al. 2022",
      "APA citations verified against original sources"]),
]
al, agap = 0.6, 0.35
aw = (12.13 - agap * 2) / 3
atop, ah = 1.85, 3.4
for i, (title, accent, fill, items) in enumerate(cols):
    x = al + i * (aw + agap)
    card(s, x, atop, aw, ah, fill=fill, border=None, radius=0.06)
    para(tb(s, x + 0.3, atop + 0.28, aw - 0.6, 0.5), title, 15, accent, bold=True, first=True)
    tf = tb(s, x + 0.3, atop + 1.0, aw - 0.55, ah - 1.2)
    for j, it in enumerate(items):
        para(tf, "•  " + it, 11, INK, first=(j == 0), sa=9, line=1.05)
card(s, 0.6, 5.55, 12.13, 1.35, fill=RGBColor(0xE7, 0xF6, 0xF0), border=TEAL, radius=0.06, lw=1.25)
tf = tb(s, 0.95, 5.78, 11.4, 1.0)
para(tf, "Human oversight policy", 13.5, TEAL, bold=True, first=True, sa=5)
para(tf, "Every AI output was reviewed and revised by the team before inclusion; statistics and citations were "
         "verified against original sources. VHN itself applies human oversight to every AI-generated result — "
         "the app routes users to real clinicians and never acts autonomously.", 11, INK, line=1.15)

# ================= SLIDE 7 — What's Next =================
s = slide(DARK)
para(tb(s, 0.7, 0.45, 12, 0.9), "What's Next", 33, WHITE, bold=True, first=True)
para(tb(s, 0.72, 1.32, 12, 0.5), "Weeks 7–9: advanced module, doctor portal, QA, and final demo",
     13.5, MUTED_LT, first=True)
nexts = [
    ("Wire OpenBioLLM-8B", "llama-cpp-python + RAG over trusted medical sources"),
    ("Doctor & clinic portal", "Patient roster, history hand-off, multi-sided routing"),
    ("Advanced features", "Voice input · PDF export · AI symptom summary"),
    ("BC Health Connect", "Live clinic queues + Health Connect Registry integration"),
    ("Safety & QA review", "Edge-case testing, privacy audit, security hardening"),
    ("Final demo & docs", "Complete technical documentation and live presentation"),
]
nl, ngap = 0.65, 0.35
nw = (12.0 - ngap * 2) / 3
ntop, nh, nvgap = 1.95, 1.85, 0.3
for i, (title, desc) in enumerate(nexts):
    col, row = i % 3, i // 3
    x = nl + col * (nw + ngap)
    y = ntop + row * (nh + nvgap)
    card(s, x, y, nw, nh, fill=DARK_CARD, border=RGBColor(0x2A, 0x36, 0x60), radius=0.06, lw=1.0)
    oval(s, x + 0.32, y + 0.32, 0.5, 0.5, TEAL)
    tf = tb(s, x + 0.32, y + 1.0, nw - 0.6, nh - 1.1)
    para(tf, title, 13.5, WHITE, bold=True, first=True, sa=5, line=1.0)
    para(tf, desc, 10.5, MUTED_LT, line=1.1)
para(tb(s, 0.7, 6.95, 12, 0.4), "VHN — guidance, not diagnosis. Every result routes to a real clinician.",
     12, PRIMARY_LT, italic=True, align=PP_ALIGN.CENTER, first=True)

out = __file__.replace("build_midterm_deck.py", "VHN_Midterm_Refined.pptx")
prs.save(out)
print("Saved:", out)
